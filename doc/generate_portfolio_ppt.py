#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Portfolio Website PPT 생성 스크립트
index.html의 내용을 기반으로 모던한 PowerPoint 프레젠테이션을 생성합니다.

사용 방법:
1. python-pptx 설치: pip install python-pptx
2. 스크립트 실행: python doc/generate_portfolio_ppt.py
   또는 doc 폴더에서: python generate_portfolio_ppt.py
3. 생성된 파일: doc/PORTFOLIO_PRESENTATION.pptx
"""

import os
import sys
from pathlib import Path

# 현재 스크립트의 디렉토리를 기준으로 경로 설정
SCRIPT_DIR = Path(__file__).parent.absolute()
DOC_DIR = SCRIPT_DIR
ROOT_DIR = SCRIPT_DIR.parent

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_portfolio_ppt():
    """포트폴리오 PPT 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 모던한 색상 팔레트
    primary_color = RGBColor(0, 51, 102)  # 진한 파란색
    secondary_color = RGBColor(70, 130, 180)  # 스틸 블루
    accent_color = RGBColor(255, 140, 0)  # 다크 오렌지
    text_color = RGBColor(51, 51, 51)  # 다크 그레이
    light_bg = RGBColor(245, 245, 250)  # 연한 배경
    
    # 슬라이드 1: 타이틀 슬라이드
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Portfolio"
    subtitle.text = "Full-Stack Developer\nJava • Spring Framework • Flutter • Mobile Development"
    
    # 타이틀 스타일 설정
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # 슬라이드 2: 소개 (About Me)
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "About Me"
    content2.text = """Full-Stack 개발자

전문 분야:
• Java, Spring Framework, Flutter 등 다양한 기술 스택을 활용한 웹 및 모바일 애플리케이션 개발
• 우리은행, 신한은행, KB국민카드 등 금융권 프로젝트 경험
• 안드로이드 네이티브 앱 개발부터 백엔드 서버 개발까지 전반적인 개발 역량

주요 통계:
✓ 15+ 프로젝트 완료
✓ 6+ 년 프리랜서 경력
✓ 4개 자격증 보유

지속적인 학습과 성장을 통해 더 나은 개발자가 되기 위해 노력하고 있으며,
최근에는 Spring Framework 기반 Java Full-Stack 개발자 양성과정을 수료하여 최신 기술을 습득했습니다."""
    
    title2.text_frame.paragraphs[0].font.size = Pt(44)
    title2.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # 슬라이드 3: 기술 스택 - Backend
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "Technical Skills - Backend"
    content3.text = """Backend Technologies:

• Java (90%) - 객체지향 프로그래밍, 멀티스레딩
• Spring Framework (90%) - MVC, Security, Data JPA
• Spring Boot (85%) - 마이크로서비스 아키텍처
• Spring AI (80%) - AI 통합 및 LLM 연동
• JSP/Servlet (85%) - 웹 애플리케이션 개발
• MyBatis (85%) - 데이터베이스 매핑
• Python (80%) - 스크립팅 및 자동화

주요 경험:
✓ RESTful API 설계 및 구현
✓ Spring Security 기반 인증/인가 시스템
✓ OAuth2 소셜 로그인 통합
✓ Spring AI를 활용한 AI 기능 구현"""
    
    title3.text_frame.paragraphs[0].font.size = Pt(36)
    title3.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # 슬라이드 4: 기술 스택 - Frontend & Mobile
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "Technical Skills - Frontend & Mobile"
    content4.text = """Frontend & Mobile Technologies:

Frontend:
• HTML5/CSS3 (90%) - 반응형 웹 디자인
• Bootstrap (85%) - UI 프레임워크
• JavaScript/jQuery (85%) - 동적 웹 개발
• Flutter/Dart (85%) - 크로스 플랫폼 개발

Mobile:
• Android/Java & Kotlin (90%) - 네이티브 앱 개발
• iOS/Swift & SwiftUI (80%) - iOS 앱 개발
• Python (80%) - 모바일 자동화

주요 경험:
✓ Flutter 기반 크로스 플랫폼 앱 개발
✓ Android 네이티브 앱 개발 (금융권 프로젝트)
✓ iOS 앱 개발 및 배포
✓ 반응형 웹 애플리케이션 개발"""
    
    title4.text_frame.paragraphs[0].font.size = Pt(36)
    title4.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # 슬라이드 5: 기술 스택 - Database & Tools
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "Technical Skills - Database & Tools"
    content5.text = """Database & DevOps Tools:

Database:
• Oracle (85%) - 엔터프라이즈 데이터베이스

DevOps & Tools:
• Git/GitHub & GitLab & Bitbucket (85%) - 버전 관리
• CI/CD (Jenkins) (75%) - 지속적 통합/배포
• Docker (80%) - 컨테이너화
• Figma (80%) - UI/UX 디자인

주요 경험:
✓ Oracle 23 AI 데이터베이스 설계 및 최적화
✓ Git 기반 협업 및 코드 리뷰
✓ Docker를 활용한 컨테이너화 및 배포
✓ CI/CD 파이프라인 구축 및 관리"""
    
    title5.text_frame.paragraphs[0].font.size = Pt(36)
    title5.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # 슬라이드 6: 주요 경력 1
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "Key Experience (1/3)"
    content6.text = """우리은행 - WON뱅킹 Re-Modeling
2022.07 - 2023.07 (12개월)

프로젝트 내용:
• 우리은행 개인비대면 채널 Re-Modeling 추진사업
• 만보기 기능 추가
• 이체기능 네이티브 → 웹 서비스 전환
• 로컬 CI/CD 환경 구축

기술 스택: Android, Java, Kotlin, WebView, CI/CD

─────────────────────────────────────

신한은행 - 땡겨요 O2O 플랫폼
2021.10 - 2022.02 (5개월)

프로젝트 내용:
• 음식주문중개 O2O 플랫폼 구축
• Pull refresh 확장기능 개발
• 땡기기 기능 구현
• WebView 설계 및 구현
• Docker를 이용한 암호화/빌드 시스템 관리

기술 스택: Android, Java, Docker, WebView"""
    
    title6.text_frame.paragraphs[0].font.size = Pt(36)
    title6.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # 슬라이드 7: 주요 경력 2
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "Key Experience (2/3)"
    content7.text = """KB 국민카드 - MyData 플랫폼
2021.04 - 2021.08 (5개월)

프로젝트 내용:
• KB 국민카드 표준API기반 MyData 플랫폼 개편 프로젝트
• 표준API기반 MyData 기능 적용
• 전체메뉴 > 메뉴검색 기능 추가

기술 스택: Android, Java, RESTful API

─────────────────────────────────────

하나은행 - Line Bank Indonesia
2020.08 - 2021.03 (8개월)

프로젝트 내용:
• 인도네시아 하나은행 Linebank 앱 개발
• MVVM 패턴 설계 및 구현
• 보안 키패드 이슈 해결
• Django & Bootstrap을 활용한 내부용 앱 배포 사이트 구축

기술 스택: Android, Java, MVVM, Django, Bootstrap"""
    
    title7.text_frame.paragraphs[0].font.size = Pt(36)
    title7.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # 슬라이드 8: 주요 경력 3
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "Key Experience (3/3)"
    content8.text = """KB국민은행 - 마이머니 App 고도화
2019.08 - 2019.11 (4개월)

프로젝트 내용:
• KB국민은행 마이머니 Android App 고도화 작업
• 안드로이드 네이티브 앱 개발
• 인트로 화면/프로그레스바 고도화
• 지문인증 솔루션 업데이트
• androidX 컨버팅

기술 스택: Android, Java, AndroidX

─────────────────────────────────────

키움증권 - 영웅문S MTS 개발
2018.05 - 2018.12 (8개월)

프로젝트 내용:
• 키움증권 영웅문S MTS 고도화 프로젝트
• 관심종목 C++ 공통 플랫폼 개발
• Javascript를 이용한 MTS 화면개발

기술 스택: C++, JavaScript, WebView"""
    
    title8.text_frame.paragraphs[0].font.size = Pt(36)
    title8.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # 슬라이드 9: 주요 프로젝트 - Miracle Reading System
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "Featured Project - Miracle Reading System"
    content9.text = """프로젝트 개요:
독서 습관 형성과 도서 관리를 위한 종합적인 웹 애플리케이션

개발 기간: 2025.11.10 - 2025.12.10 (1개월)
개발 형태: 개인 프로젝트 (1인 총괄 개발)

주요 기능:
✓ 사용자 인증 및 관리 (폼 로그인, OAuth2)
✓ 도서 관리 시스템 (알라딘 API 연동)
✓ AI 기반 도서 요약 (Spring AI + Ollama)
✓ 독서 계획 및 기록 관리
✓ 속독 훈련 기능
✓ 갤러리 및 소셜 기능 (좋아요, 댓글)
✓ 마인드맵 기능
✓ 관리자 콘솔

기술 스택:
Java 17, Spring Boot 3.3.5, Spring AI, Oracle 23 AI,
JSP, Bootstrap 5, jQuery, Docker, Ollama (Qwen3:1.7b)

주요 성과:
• AI 통합: Spring AI를 활용한 로컬 LLM 연동
• 확장 가능한 아키텍처: 계층형 구조 설계
• 다중 인증 시스템: 폼 로그인 + OAuth2 통합"""
    
    title9.text_frame.paragraphs[0].font.size = Pt(32)
    title9.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # 슬라이드 10: 주요 프로젝트 - Productivity Hub
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "Featured Project - Productivity Hub"
    content10.text = """프로젝트 개요:
Flutter 기반의 통합 생산성 앱

개발 기간: 2025.12.04 오후 (4시간)
개발 형태: 개인 프로젝트 (1인 총괄 개발)

주요 기능:
✓ 할 일 관리 (Todo) - 추가/수정/삭제, 완료 상태 토글
✓ 아이디어 기록 - 카테고리별 아이디어 관리
✓ 독서 카드 - 독서 진행 관리, 키워드/요약 기록
✓ 날씨 정보 - 현재 위치 및 도시별 날씨 조회
✓ 뉴스 피드 - AI/양자컴퓨팅 관련 최신 뉴스

기술 스택:
Flutter 3.x, Dart, Provider, SQLite, Open-Meteo API,
RSS Feed, Geolocator

아키텍처:
• Provider 패턴 (MVVM 기반) 상태 관리
• SQLite 로컬 데이터 저장
• RESTful API 연동 (날씨, 뉴스)
• 반응형 UI 디자인"""
    
    title10.text_frame.paragraphs[0].font.size = Pt(32)
    title10.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # 슬라이드 11: 교육 및 자격증
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title11 = slide11.shapes.title
    content11 = slide11.placeholders[1]
    
    title11.text = "Education & Certifications"
    content11.text = """학위:

• 석사 - 광주과학기술원 기전공학과 (2005.03 ~ 2007.08)
• 학사 - 강원대학교 전기전자공학과 (1995.03 ~ 2004.02)

교육 이력:

• Spring Framework 기반 Java Full-Stack 개발자 양성과정
  쌍용강북교육센터 (2025.05.12 - 2025.11.12, 944시간)
  
• 소음진동평가모니터링시스템개발 과정
  경영기술개발원교육센터 (2012.06 - 2012.12, 960시간)
  
• 임베디드 SW 전문가 과정
  한국정보기술연구원 (2007.10 - 2008.03, 960시간)

자격증:
✓ 정보처리기사 (2025.09)
✓ RFID-GL (2013.11)
✓ SCJP (2010.04)
✓ 전기공사 (2004.08)

해외 경험:
• 산업인력공단 월드잡 연수 프로그램 (2010.07 - 2011.05)
  - Canadagate IT 비즈니스 실무 과정
  - Advanced 과정(Toefl) 수업 약 4개월 수강
  - 미국, 캐나다 Brain-based Speed Reading 세미나 참석"""
    
    title11.text_frame.paragraphs[0].font.size = Pt(36)
    title11.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # 슬라이드 12: 핵심 역량 요약
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    title12 = slide12.shapes.title
    content12 = slide12.placeholders[1]
    
    title12.text = "Core Competencies"
    content12.text = """기술 역량:

Full-Stack Development:
• Backend: Java, Spring Framework, Spring Boot, Spring AI
• Frontend: HTML5/CSS3, JavaScript, jQuery, Bootstrap
• Mobile: Android (Java/Kotlin), iOS (Swift/SwiftUI), Flutter
• Database: Oracle, SQLite
• DevOps: Git, Docker, CI/CD (Jenkins)

프로젝트 경험:

금융권 프로젝트:
• 우리은행, 신한은행, KB국민카드/은행, 하나은행 등
• 안드로이드 네이티브 앱 개발
• 웹 서비스 전환 및 고도화
• 보안 및 인증 시스템 구현

기타 프로젝트:
• O2O 플랫폼 개발
• 증권사 MTS 개발
• 도시가스 검침 시스템 개발
• 통합 생산성 앱 개발

주요 강점:
✓ 금융권 프로젝트 다수 경험
✓ 풀스택 개발 역량
✓ 크로스 플랫폼 개발 경험
✓ 최신 기술 학습 및 적용 능력"""
    
    title12.text_frame.paragraphs[0].font.size = Pt(36)
    title12.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # 슬라이드 13: 마무리
    slide13 = prs.slides.add_slide(prs.slide_layouts[0])
    title13 = slide13.shapes.title
    subtitle13 = slide13.placeholders[1]
    
    title13.text = "Thank You"
    subtitle13.text = "Portfolio\n\nFull-Stack Developer\n\n문의사항이 있으시면 언제든지 연락주세요.\n\nLinkedIn: linkedin.com/in/namil-kim-a59951123\nGitHub: github.com/NAM-IL"
    
    title13.text_frame.paragraphs[0].font.size = Pt(54)
    title13.text_frame.paragraphs[0].font.bold = True
    title13.text_frame.paragraphs[0].font.color.rgb = primary_color
    subtitle13.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle13.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # 파일 저장 (doc 폴더에 저장)
    filename = DOC_DIR / "PORTFOLIO_PRESENTATION.pptx"
    prs.save(str(filename))
    print(f"✅ 포트폴리오 PPT가 생성되었습니다: {filename}")
    print(f"📊 총 {len(prs.slides)}개의 슬라이드가 포함되어 있습니다.")
    print(f"📝 생성된 슬라이드 목록:")
    for i, slide in enumerate(prs.slides, 1):
        try:
            title = slide.shapes.title.text if slide.shapes.title else "제목 없음"
            print(f"   {i}. {title}")
        except:
            print(f"   {i}. (슬라이드 {i})")
    print(f"\n💡 index.html의 내용을 기반으로 작성되었습니다.")
    print(f"📁 저장 위치: {filename}")

if __name__ == "__main__":
    try:
        create_portfolio_ppt()
    except ImportError:
        print("❌ python-pptx 라이브러리가 설치되지 않았습니다.")
        print("📦 설치 방법: pip install python-pptx")
    except Exception as e:
        print(f"❌ 오류 발생: {e}")
        import traceback
        traceback.print_exc()
